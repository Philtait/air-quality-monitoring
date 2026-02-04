"""
Air Quality Analysis - 15 Slide PowerPoint Presentation Generator

This script creates a professional PowerPoint presentation based on
the air quality data visualizations and analysis findings.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(0, 63, 92)
TEAL = RGBColor(0, 128, 128)
ORANGE = RGBColor(255, 127, 14)
GREEN = RGBColor(44, 160, 44)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)

# Path to visualizations
viz_path = r"c:\Users\phili\air-quality-project\visualizations"

def add_title_slide(title, subtitle=""):
    """Add a title slide with dark blue background"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        # Add subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, image_path=None, bullet_points=None, two_column=False):
    """Add a content slide with title, optional image and bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
        prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()
    
    # Add title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if image_path and os.path.exists(image_path):
        if bullet_points and two_column:
            # Two column layout: bullets on left, image on right
            # Add bullet points on left
            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(5.5))
            tf = text_box.text_frame
            tf.word_wrap = True
            for i, point in enumerate(bullet_points):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = "• " + point
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(50, 50, 50)
                p.space_before = Pt(12)
            
            # Add image on right
            slide.shapes.add_picture(image_path, Inches(5.8), Inches(1.4), width=Inches(7.2))
        else:
            # Full width image
            slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), width=Inches(12.333))
    elif bullet_points:
        # Bullets only
        text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.5))
        tf = text_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_before = Pt(16)
    
    return slide

def add_key_findings_slide(title, findings):
    """Add a key findings slide with icons and text"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
        prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()
    
    # Add title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Add findings in a grid layout
    cols = 2
    rows = (len(findings) + 1) // 2
    box_width = Inches(5.8)
    box_height = Inches(1.3)
    
    for i, finding in enumerate(findings):
        col = i % cols
        row = i // cols
        left = Inches(0.5) + col * (box_width + Inches(0.5))
        top = Inches(1.6) + row * (box_height + Inches(0.3))
        
        # Add finding box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 248, 250)
        box.line.color.rgb = RGBColor(200, 210, 220)
        
        # Add finding text
        text_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), 
                                            box_width - Inches(0.4), box_height - Inches(0.4))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = finding
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(40, 40, 40)
    
    return slide


# ============== CREATE THE 15 SLIDES ==============

# Slide 1: Title Slide
add_title_slide(
    "Air Quality in America",
    "A Comprehensive Analysis of 42 Years of Data (1980-2022)"
)

# Slide 2: Executive Summary / Key Findings
add_key_findings_slide(
    "Executive Summary: Key Findings",
    [
        "📈 26% improvement in air quality from 1980 to 2022",
        "☀️ Summer months have the worst air quality (July peak: 57.5 AQI)",
        "🏭 Ozone is the dominant pollutant (54% of readings)",
        "📍 California cities consistently rank as worst for air quality",
        "🌍 Geography matters more than population size",
        "⚠️ Extreme events have decreased dramatically since the 1990s"
    ]
)

# Slide 3: About the Dataset
add_content_slide(
    "About the Dataset",
    bullet_points=[
        "5 million+ air quality readings from EPA monitoring stations",
        "42 years of data spanning 1980 to 2022",
        "Coverage: 634 cities across 52 states/territories",
        "5 pollutants tracked: Ozone, PM2.5, PM10, NO2, CO",
        "Air Quality Index (AQI) used as primary metric",
        "AQI Categories: Good (0-50), Moderate (51-100), Unhealthy (101+)"
    ]
)

# Slide 4: 42-Year Time Series Trend
add_content_slide(
    "42 Years of Progress: Air Quality Improvement Over Time",
    image_path=os.path.join(viz_path, "viz1_timeseries.png"),
    bullet_points=[
        "Average AQI decreased from 54.6 to 40.6",
        "26% overall improvement in air quality",
        "Clear downward trend since 1990s",
        "Clean Air Act regulations showing impact"
    ],
    two_column=True
)

# Slide 5: Decade by Decade Analysis
add_content_slide(
    "Decade-by-Decade Analysis: Consistent Improvement",
    image_path=os.path.join(viz_path, "viz3_decades.png"),
    bullet_points=[
        "1980s: Average AQI 53.4 (worst decade)",
        "1990s: AQI dropped to 46.8",
        "2000s: Slight increase to 48.0",
        "2010s: Significant drop to 43.1",
        "2020s: Best yet at 40.9 (23% better than 1980s)"
    ],
    two_column=True
)

# Slide 6: Seasonal Patterns
add_content_slide(
    "Seasonal Patterns: Summer Has Worst Air Quality",
    image_path=os.path.join(viz_path, "viz2_seasonal.png"),
    bullet_points=[
        "July peak: 57.5 AQI (highest month)",
        "Winter months: Lowest AQI readings",
        "Ozone dominates in summer",
        "PM2.5 more prevalent in winter",
        "Temperature and sunlight drive patterns"
    ],
    two_column=True
)

# Slide 7: Pollutant Analysis
add_content_slide(
    "Pollutant Overview: Ozone Dominates",
    image_path=os.path.join(viz_path, "viz8_pollutant_overview.png"),
    bullet_points=[
        "Ozone: 54% of all readings (most common)",
        "PM2.5: 25% of readings",
        "Ozone has highest average AQI (51.6)",
        "PM10 has extreme outliers (max: 20,646!)",
        "CO shows lowest average levels"
    ],
    two_column=True
)

# Slide 8: Pollutant Seasonality
add_content_slide(
    "Pollutant Seasonality: Different Pollutants Peak at Different Times",
    image_path=os.path.join(viz_path, "viz9_pollutant_seasonality.png"),
    bullet_points=[
        "Ozone surges dramatically in summer",
        "PM2.5 peaks in winter months",
        "NO2 relatively stable year-round",
        "PM10 shows spring peaks (dust seasons)",
        "Understanding seasonality helps planning"
    ],
    two_column=True
)

# Slide 9: State Rankings
add_content_slide(
    "State Rankings: Who Has the Worst Air Quality?",
    image_path=os.path.join(viz_path, "viz4_state_heatmap.png"),
    bullet_points=[
        "District of Columbia leads (70.3 AQI)",
        "California second worst (63.8 AQI)",
        "Utah, Rhode Island, Maryland follow",
        "Mississippi has best air (46.4 AQI)",
        "Most states fall in 'Good' category"
    ],
    two_column=True
)

# Slide 10: City Rankings
add_content_slide(
    "City Rankings: Best and Worst Cities in America",
    image_path=os.path.join(viz_path, "viz5_city_rankings.png"),
    bullet_points=[
        "TOP 10 WORST CITIES - All in California!",
        "Riverside #1 worst (124.9 AQI)",
        "Los Angeles, Bakersfield, Fresno follow",
        "TOP 10 BEST CITIES - Diverse locations",
        "Alma, MI cleanest (5.0 AQI)",
        "Clean cities found across many states"
    ],
    two_column=True
)

# Slide 11: Regional Analysis
add_content_slide(
    "Regional Analysis: Geography Matters",
    image_path=os.path.join(viz_path, "viz6_regional_comparison.png"),
    bullet_points=[
        "Mountain West: Worst region (50.6 AQI)",
        "Great Plains: Best region (42.7 AQI)",
        "West Coast high due to PM10 (dust)",
        "Midwest/South: Ozone dominant",
        "Regional variations in pollutant types"
    ],
    two_column=True
)

# Slide 12: Geographic Distribution Map
add_content_slide(
    "Geographic Distribution: Mapping Air Quality Across the USA",
    image_path=os.path.join(viz_path, "viz7_geographic_scatter.png"),
    bullet_points=[
        "All 5 worst cities located in California",
        "Midwest shows consistently clean air",
        "East Coast: Moderate levels",
        "Pacific Northwest: Generally good",
        "Clear geographic clustering patterns"
    ],
    two_column=True
)

# Slide 13: Population vs Air Quality
add_content_slide(
    "Surprising Finding: Population Has WEAK Effect on Air Quality",
    image_path=os.path.join(viz_path, "viz11_population_relationships.png"),
    bullet_points=[
        "Correlation r=0.20 (weak positive)",
        "Population density r=0.14 (even weaker!)",
        "Geography and climate matter more",
        "Worst cities are NOT always the biggest",
        "Riverside worse than NYC despite size"
    ],
    two_column=True
)

# Slide 14: Extreme Events
add_content_slide(
    "Extreme Air Quality Events: A Success Story",
    image_path=os.path.join(viz_path, "viz10_extreme_events.png"),
    bullet_points=[
        "1988: Worst year with 5,029 extreme events",
        "Dramatic decline since the 1990s",
        "Ozone causes 56% of extreme events",
        "PM2.5 responsible for 38% of extremes",
        "2020s: Fewer than 500 events per year"
    ],
    two_column=True
)

# Slide 15: Conclusions and Recommendations
add_key_findings_slide(
    "Conclusions & Future Recommendations",
    [
        "✅ Air quality regulations (Clean Air Act) are working",
        "🎯 Focus on California - consistently worst performer",
        "☀️ Summer interventions crucial (Ozone peaks)",
        "🏙️ City size alone doesn't predict air quality",
        "📊 Continue monitoring PM2.5 and Ozone as top concerns",
        "🔮 Data-driven decisions can further improve outcomes"
    ]
)

# Save the presentation
output_path = r"c:\Users\phili\air-quality-project\Air_Quality_Analysis_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation created successfully!")
print(f"📁 Saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
